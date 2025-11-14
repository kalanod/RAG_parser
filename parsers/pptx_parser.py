from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import List

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from PIL import Image
import pytesseract
import logging

from parsers.abstract_document_parser import AbstractDocumentParser

LOGGER = logging.getLogger(__name__)


# ------------------------ CONFIG ------------------------
@dataclass
class PPTXParserConfig:
    chunk_size: int = 700
    chunk_overlap: int = 150
    ocr_lang: str = "rus+eng"
    min_ocr_w: int = 100
    min_ocr_h: int = 50


# ------------------------ PARSER ------------------------
class PPTXParser(AbstractDocumentParser):
    supported_extensions = (".pptx",)
    config = PPTXParserConfig()

    def parse(self, path: Path, **kwargs: object) -> List[Document]:
        try:
            documents: List[Document] = []
            documents.extend(self._extract_slide_text_blocks(path))
            documents.extend(self._extract_table_blocks(path))
            documents.extend(self._extract_ocr_blocks(path))
            return self.split_to_docs(documents)
        except Exception as exc:
            LOGGER.exception("Failed to parse PPTX %s: %s", path, exc)
            return [
                Document(
                    page_content=f"[ERROR] Failed to load PPTX: {str(exc)}",
                    metadata={
                        "source": str(path),
                        "error": True,
                        "error_type": type(exc).__name__,
                    },
                )
            ]

    # ------------------------ SPLIT ------------------------
    def split_to_docs(self, documents: List[Document]) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.chunk_size,
            chunk_overlap=self.config.chunk_overlap,
            separators=["\n\n", "\n", ".", "?", "!", " "],
        )
        split_docs: List[Document] = []

        for doc in documents:
            chunks = splitter.split_text(doc.page_content)
            for i, chunk in enumerate(chunks):
                split_docs.append(
                    Document(
                        page_content=chunk,
                        metadata={
                            **doc.metadata,
                            "chunk_index": i,
                            "chunk_count": len(chunks),
                        },
                    )
                )
        return split_docs

    def _iter_shapes_recursive(self, shapes):
        for shp in shapes:
            yield shp
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub in self._iter_shapes_recursive(shp.shapes):
                        yield sub
            except Exception:
                continue

    def _extract_slide_text_blocks(self, path: Path) -> List[Document]:
        prs = Presentation(str(path))
        docs: List[Document] = []

        for slide_idx, slide in enumerate(prs.slides):
            title_text = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    title_text = slide.shapes.title.text.strip()
            except Exception:
                title_text = ""

            body_lines: List[str] = []
            for shp in self._iter_shapes_recursive(slide.shapes):
                if not getattr(shp, "has_text_frame", False):
                    continue
                text_frame = shp.text_frame
                if not text_frame:
                    continue
                for p in text_frame.paragraphs:
                    txt = (p.text or "").strip()
                    if txt:
                        body_lines.append(txt)

            content_parts: List[str] = []
            if title_text:
                content_parts.append(f"# {title_text}")
            if body_lines:
                content_parts.append("\n".join(body_lines))

            slide_text = "\n\n".join([p for p in content_parts if p.strip()])

            if slide_text.strip():
                docs.append(
                    Document(
                        page_content=slide_text.strip(),
                        metadata={
                            "source": str(path),
                            "type": "slide_text",
                            "slide_index": slide_idx,
                            "title": title_text,
                        },
                    )
                )

        return docs

    def _extract_table_blocks(self, path: Path) -> List[Document]:
        prs = Presentation(str(path))
        docs: List[Document] = []

        for slide_idx, slide in enumerate(prs.slides):
            table_count = 0
            for shp in self._iter_shapes_recursive(slide.shapes):
                if not getattr(shp, "has_table", False):
                    continue

                table = shp.table
                if not table:
                    continue

                rows_text: List[List[str]] = []
                for row in table.rows:
                    row_cells: List[str] = []
                    for cell in row.cells:
                        txt = (cell.text or "").strip()
                        row_cells.append(txt)
                    rows_text.append(row_cells)

                if not any(any(c.strip() for c in r) for r in rows_text):
                    continue
                header = rows_text[0]
                num_cols = len(header)

                def to_md_row(cols: List[str]) -> str:
                    return "| " + " | ".join((c or "").replace("\n", " ") for c in cols) + " |"

                header_line = to_md_row(header)
                sep_line = "| " + " | ".join("---" for _ in range(num_cols)) + " |"
                body_lines = [to_md_row(r) for r in rows_text[1:]]

                md_table = "\n".join([header_line, sep_line, *body_lines]).strip()

                if md_table:
                    docs.append(
                        Document(
                            page_content=md_table,
                            metadata={
                                "source": str(path),
                                "type": "table",
                                "slide_index": slide_idx,
                                "table_index": table_count,
                            },
                        )
                    )
                    table_count += 1

        return docs

    def _extract_ocr_blocks(self, path: Path) -> List[Document]:
        prs = Presentation(str(path))
        docs: List[Document] = []

        for slide_idx, slide in enumerate(prs.slides):
            for shp in self._iter_shapes_recursive(slide.shapes):
                try:
                    if shp.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                except Exception:
                    continue

                try:
                    image = shp.image
                    blob = image.blob
                except Exception as e:
                    LOGGER.debug("Skipping image on slide %s due to error: %s", slide_idx, e)
                    continue

                try:
                    img = Image.open(BytesIO(blob))
                    w, h = img.size
                except Exception as e:
                    LOGGER.debug("Failed to open image on slide %s: %s", slide_idx, e)
                    continue

                if w < self.config.min_ocr_w or h < self.config.min_ocr_h:
                    continue

                try:
                    ocr_text = pytesseract.image_to_string(img, lang=self.config.ocr_lang).strip()
                except Exception as e:
                    LOGGER.warning("OCR failed for PPTX image on slide %s in %s: %s", slide_idx, path, e)
                    ocr_text = ""

                if ocr_text:
                    docs.append(
                        Document(
                            page_content=ocr_text,
                            metadata={
                                "source": str(path),
                                "type": "image_text",
                                "slide_index": slide_idx,
                                "image_size": f"{w}x{h}",
                            },
                        )
                    )
        return docs
